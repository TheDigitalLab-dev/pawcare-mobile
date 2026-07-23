#!/usr/bin/env python3
"""Genera docs/pst/presentacion/diapositivas-pawcare-mobile.pptx a partir del
deck HTML (diapositivas-pawcare-mobile.html), para importarlo en Google Slides.

Pipeline: Chrome headless imprime el HTML (CSS @page 960x540 pt, 16:9) a un PDF
temporal, pdftoppm rasteriza cada diapositiva a PNG (200 dpi) y python-pptx las
inserta a pagina completa en un .pptx 16:9. Requiere: google-chrome, pdftoppm,
python-pptx.
"""
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "docs" / "pst" / "presentacion" / "diapositivas-pawcare-mobile.html"
OUT = ROOT / "docs" / "pst" / "presentacion" / "diapositivas-pawcare-mobile.pptx"

DPI = 200
SLIDE_W = Pt(960)
SLIDE_H = Pt(540)


def main():
    with tempfile.TemporaryDirectory() as tmp:
        pdf = Path(tmp) / "deck.pdf"
        subprocess.run(
            [
                "google-chrome", "--headless=new", "--disable-gpu",
                "--no-pdf-header-footer", f"--print-to-pdf={pdf}",
                SRC.resolve().as_uri(),
            ],
            check=True,
            capture_output=True,
        )
        prefix = Path(tmp) / "slide"
        subprocess.run(
            ["pdftoppm", "-png", "-r", str(DPI), str(pdf), str(prefix)],
            check=True,
        )
        pages = sorted(Path(tmp).glob("slide-*.png"))
        if not pages:
            raise RuntimeError("pdftoppm no produjo diapositivas")

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        blank = prs.slide_layouts[6]
        for png in pages:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(png), 0, 0, width=SLIDE_W, height=SLIDE_H)
        prs.save(str(OUT))
        print(f"OK -> {OUT} ({len(pages)} diapositivas)")


if __name__ == "__main__":
    main()
